import streamlit as st
import os
import openai
import fitz  # PyMuPDF
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
import docx2txt
from sentence_transformers import SentenceTransformer
import faiss
import pickle
import numpy as np
import logging
from datetime import datetime
import signal
import sys

# Set up logging
logging.basicConfig(level=logging.INFO)

# Secrets Manager to retrieve keys and endpoints
class SecretsManager:
    @staticmethod
    def get_secret(secret_name):
        secrets = {
            "AZURE_API_KEY": "YOUR_AZURE_API_KEY",
            "AZURE_API_DEPLOYMENT": "YOUR_AZURE_API_DEPLOYMENT",
            "AZURE_API_ENDPOINT": "YOUR_AZURE_API_ENDPOINT"
        }
        return secrets.get(secret_name)

# Retrieve API key, deployment ID, and endpoint from Secrets Manager
key = SecretsManager.get_secret('AZURE_API_KEY')
deployment = SecretsManager.get_secret("AZURE_API_DEPLOYMENT")
endpoint = SecretsManager.get_secret("AZURE_API_ENDPOINT")

# Azure OpenAI Configuration
openai.api_type = "azure"
openai.api_key = key
openai.api_base = endpoint
openai.api_version = "2024-02-01"

# Define the deployment ID for your Azure OpenAI model
deployment_id = deployment

# Initialize Sentence Transformer Model
model = SentenceTransformer('all-MiniLM-L6-v2')

# Create necessary directories
os.makedirs('documents', exist_ok=True)
os.makedirs('vector_store', exist_ok=True)

# Load existing FAISS index or create a new one
if os.path.exists('vector_store/faiss_index'):
    with open('vector_store/faiss_index', 'rb') as f:
        index = pickle.load(f)
else:
    index = faiss.IndexIDMap(faiss.IndexFlatL2(384))  # 384-dimensional embeddings with IDMap

# Load existing document metadata or create a new one
metadata_path = 'vector_store/doc_metadata.pkl'
if os.path.exists(metadata_path):
    with open(metadata_path, 'rb') as f:
        doc_metadata = pickle.load(f)
else:
    doc_metadata = {}

# Function to extract text from different document types
def extract_text_from_doc(file_path):
    try:
        if file_path.endswith('.pdf'):
            text = ""
            pdf = fitz.open(file_path)
            page = pdf.load_page(0)  # Load only the first page
            text += page.get_text()
            return text
        elif file_path.endswith('.pptx'):
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
            return text
        elif file_path.endswith('.docx'):
            return docx2txt.process(file_path)
        elif file_path.endswith('.txt'):
            with open(file_path, 'r') as file:
                return file.read()
        else:
            return ""
    except Exception as e:
        logging.error(f"Error extracting text from {file_path}: {e}")
        return ""

# Function to embed documents and update FAISS index and metadata
def embed_documents(file_paths):
    def split_text(text, max_length=512):
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = ""
        for para in paragraphs:
            if len(current_chunk) + len(para) < max_length:
                current_chunk += para + "\n\n"
            else:
                chunks.append(current_chunk.strip())
                current_chunk = para + "\n\n"
        if current_chunk:
            chunks.append(current_chunk.strip())
        return chunks

    try:
        for file_path in file_paths:
            text = extract_text_from_doc(file_path)
            if text:
                chunks = split_text(text)
                embeddings = model.encode(chunks)
                ids = np.arange(index.ntotal, index.ntotal + len(embeddings))
                index.add_with_ids(np.array(embeddings, dtype=np.float32), ids)
                for i, chunk in enumerate(chunks):
                    doc_metadata[int(ids[i])] = {
                        'text': chunk,
                        'file_path': file_path
                    }
                with open('vector_store/faiss_index', 'wb') as f:
                    pickle.dump(index, f)
                with open(metadata_path, 'wb') as f:
                    pickle.dump(doc_metadata, f)
            else:
                logging.warning(f"No text extracted from {file_path}. Skipping embedding.")
    except Exception as e:
        logging.error(f"Error embedding documents: {e}")

# Function to handle document uploads
def handle_upload(uploaded_files):
    try:
        file_paths = []
        for uploaded_file in uploaded_files:
            file_path = os.path.join('documents', uploaded_file.name)
            with open(file_path, 'wb') as file:
                file.write(uploaded_file.getbuffer())
            file_paths.append(file_path)
        
        embed_documents(file_paths)
        return file_paths
    except Exception as e:
        logging.error(f"Error handling upload: {e}")
        return []

# AzureChatOpenAI class definition (assuming you have it implemented elsewhere)
class AzureChatOpenAI:
    def __init__(self, model_name, temperature, openai_api_key, azure_deployment, azure_endpoint, openai_api_version):
        self.model_name = model_name
        self.temperature = temperature
        self.openai_api_key = openai_api_key
        self.azure_deployment = azure_deployment
        self.azure_endpoint = azure_endpoint
        self.openai_api_version = openai_api_version

    def generate_response(self, prompt):
        response = openai.ChatCompletion.create(
            deployment_id=self.azure_deployment,
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": prompt}
            ]
        )
        return response['choices'][0]['message']['content'].strip()

# Initialize the AzureChatOpenAI model
chat_model = AzureChatOpenAI(
    model_name='gpt-35-turbo',
    temperature=0.05,
    openai_api_key=key,
    azure_deployment=deployment,
    azure_endpoint=endpoint,
    openai_api_version='2024-02-01'
)

# Function to retrieve previous chat history (if implemented)
def get_previous_chat_history():
    return st.session_state.get("chat_history", [])

# Function to save chat history (if implemented)
def save_chat_history(question, answer, references=None):
    if "chat_history" not in st.session_state:
        st.session_state["chat_history"] = []
    st.session_state["chat_history"].append({"question": question, "answer": answer, "references": references})

def get_chatbot_response(user_input):
    try:
        # Retrieve relevant documents
        user_embedding = model.encode([user_input])
        D, I = index.search(np.array(user_embedding, dtype=np.float32), 5)
        relevant_docs = []
        logging.info(f"Indices: {I[0]}")
        seen_documents = set()
        for idx in I[0]:
            if idx != -1:  # Ensure index is valid
                doc_data = doc_metadata.get(int(idx), None)
                if doc_data and doc_data['file_path'] not in seen_documents:
                    relevant_docs.append(doc_data)
                    seen_documents.add(doc_data['file_path'])

        logging.info(f"Relevant documents: {relevant_docs}")

        # Summarize or filter the relevant documents if they are too long
        max_tokens = 1000  # Set a token limit for the prompt
        prompt_docs = []
        current_tokens = 0
        doc_references = []
        for doc_data in relevant_docs:
            text = extract_text_from_doc(doc_data['file_path'])
            tokens = len(text.split())
            if current_tokens + tokens > max_tokens:
                break
            prompt_docs.append(text)
            current_tokens += tokens
            doc_references.append((os.path.basename(doc_data['file_path']), doc_data['file_path']))

        # Construct the prompt with chat history
        chat_history = get_previous_chat_history()
        history_text = "\n".join([f"User: {entry['question']}\nAssistant: {entry['answer']}" for entry in chat_history])
        prompt = f"{history_text}\nUser asked: {user_input}\nRelevant documents: {' '.join(prompt_docs)}\nAnswer:"
        response = chat_model.generate_response(prompt)

        # Save chat history
        save_chat_history(user_input, response, doc_references)

        return response, doc_references
    except Exception as e:
        logging.error(f"Error getting chatbot response: {e}")
        return "I don't know.", []

def terminate_program():
    st.write("Shutting down...")
    # Send a termination signal to the program
    os.kill(os.getpid(), signal.SIGTERM)

# Streamlit UI
def streamlit_ui():
    st.title("RAG Model with Azure OpenAI and Local Documents")

    st.write("Chat History:")
    chat_history = get_previous_chat_history()
    if chat_history:
        for entry in chat_history:
            st.markdown(
                f'<div style="background-color: #000000; padding: 10px; border-radius: 5px;"><b>User:</b> {entry["question"]}</div>',
                unsafe_allow_html=True
            )
            st.markdown(
                f'<div style="background-color: #000000; padding: 10px; border-radius: 5px;"><b>Assistant:</b> {entry["answer"]}</div>',
                unsafe_allow_html=True
            )
            st.markdown(
                f'<div style="background-color: #000000; padding: 10px; border-radius: 5px;"><b>References:</b> {", ".join([ref[0] for ref in entry["references"]])}</div>',
                unsafe_allow_html=True
            )
    
    st.write("Upload documents for the RAG model.")
    uploaded_files = st.file_uploader("Choose files", type=['pdf', 'pptx', 'docx', 'txt'], accept_multiple_files=True)
    if uploaded_files:
        with st.spinner('Uploading and processing...'):
            file_paths = handle_upload(uploaded_files)
            if file_paths:
                for file_path in file_paths:
                    st.success(f"Uploaded {os.path.basename(file_path)}")
            else:
                st.error(f"Failed to upload files. Check logs for details.")
    
    st.write("Start chatting with the RAG model.")
    user_input = st.text_input("You: ")
    if st.button("Send"):
        if user_input:
            with st.spinner('Getting response...'):
                response, references = get_chatbot_response(user_input)
                st.markdown(
                    f'<div style="background-color: #000000; padding: 10px; border-radius: 5px;"><b>Assistant:</b> {response}</div>',
                    unsafe_allow_html=True
                )
                if references:
                    st.markdown(
                        f'<div style="background-color: #000000; padding: 10px; border-radius: 5px;"><b>References:</b> {", ".join([ref[0] for ref in references])}</div>',
                        unsafe_allow_html=True
                    )
                else:
                    st.markdown(
                        f'<div style="background-color: #000000; padding: 10px; border-radius: 5px;"><b>References:</b> None</div>',
                        unsafe_allow_html=True
                    )
        else:
            st.warning("Please enter a message to send.")
    
    if st.button("Stop Chat"):
        terminate_program()

    

if __name__ == '__main__':
    streamlit_ui()
